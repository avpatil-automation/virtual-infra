from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def set_font(para, font_name='Arial', font_size=Pt(18), color=RGBColor(0, 0, 0)):
    run = para.add_run()
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = color
    return run

def create_professional_presentation():
    # Create a new presentation
    ppt_pres = Presentation()

    # Slide 1: Title Slide
    slide = ppt_pres.slides.add_slide(ppt_pres.slide_layouts[0])  # Title Slide layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Future of Telecommunications"
    subtitle.text = "Insights and Perspectives\nYour Name\nDate"

    # Slide 2: Introduction
    slide = ppt_pres.slides.add_slide(ppt_pres.slide_layouts[1])  # Title and Content layout
    slide.shapes.title.text = "Introduction"
    content = slide.placeholders[1]
    content.text = ""  # Clear default text
    p = content.text_frame.add_paragraph()
    set_font(p, font_size=Pt(20))
    p.text = (
        "The telecommunications industry is at a pivotal point, driven by rapid technological advancements and changing consumer expectations. "
        "This presentation explores future trends, their implications for projects and operations, and strategies for enhancing client partnerships."
    )

    # Slide 3: Future Trends in Telecommunications
    slide = ppt_pres.slides.add_slide(ppt_pres.slide_layouts[1])
    slide.shapes.title.text = "Future Trends in Telecommunications"
    content = slide.placeholders[1]
    content.text = ""  # Clear default text
    p = content.text_frame.add_paragraph()
    set_font(p, font_size=Pt(20))
    p.text = "Key Trends:"
    
    trends = [
        "5G and Beyond: Enhanced speed and connectivity will revolutionize services.",
        "Internet of Things (IoT): Growth in connected devices will lead to new service models.",
        "Edge Computing: Processing data closer to the source reduces latency.",
        "Artificial Intelligence (AI): AI will optimize network management and customer service."
    ]
    
    for trend in trends:
        p = content.text_frame.add_paragraph()
        set_font(p, font_size=Pt(18))
        p.text = f"- {trend}"

    # Slide 4: Impact of Future Trends on Projects
    slide = ppt_pres.slides.add_slide(ppt_pres.slide_layouts[1])
    slide.shapes.title.text = "Impact of Future Trends on Projects"
    content = slide.placeholders[1]
    content.text = ""  # Clear default text
    p = content.text_frame.add_paragraph()
    set_font(p, font_size=Pt(20))
    p.text = "As these trends evolve, they will influence:"
    
    impacts = [
        "Project Methodologies: Agile and DevOps approaches will be essential to adapt quickly.",
        "Client Expectations: Clients will demand faster, more reliable services with innovative solutions.",
        "Operational Efficiency: Automation and AI will streamline processes, reducing costs and improving service delivery.",
        "Risk Management: New technologies introduce risks that require robust management strategies."
    ]
    
    for impact in impacts:
        p = content.text_frame.add_paragraph()
        set_font(p, font_size=Pt(18))
        p.text = f"- {impact}"

    # Slide 5: Influence on Accenture Operations
    slide = ppt_pres.slides.add_slide(ppt_pres.slide_layouts[1])
    slide.shapes.title.text = "Influence on Accenture Operations"
    content = slide.placeholders[1]
    content.text = ""  # Clear default text
    p = content.text_frame.add_paragraph()
    set_font(p, font_size=Pt(20))
    p.text = "The future trends will necessitate changes in Accenture's operations:"
    
    influences = [
        "Service Offerings: Development of new services around 5G, IoT, and AI.",
        "Client Engagement: Enhanced collaboration with clients through co-innovation.",
        "Talent Acquisition: Need for skilled professionals in AI, data analytics, and cybersecurity.",
        "Sustainability Initiatives: Focus on reducing the carbon footprint of telecom operations."
    ]
    
    for influence in influences:
        p = content.text_frame.add_paragraph()
        set_font(p, font_size=Pt(18))
        p.text = f"- {influence}"

    # Slide 6: Partnering with Clients
    slide = ppt_pres.slides.add_slide(ppt_pres.slide_layouts[1])
    slide.shapes.title.text = "Partnering with Clients"
    content = slide.placeholders[1]
    content.text = ""  # Clear default text
    p = content.text_frame.add_paragraph()
    set_font(p, font_size=Pt(20))
    p.text = "To enhance partnerships and increase value propositions:"
    
    strategies = [
        "Leverage Insights: Use data analytics to provide personalized solutions.",
        "Innovative Solutions: Collaborate on developing new products and services that address client needs.",
        "Future-Ready Strategies: Prepare clients for upcoming trends through education and resources.",
        "Feedback Loops: Establish mechanisms for continuous client feedback to improve offerings."
    ]
    
    for strategy in strategies:
        p = content.text_frame.add_paragraph()
        set_font(p, font_size=Pt(18))
        p.text = f"- {strategy}"

    # Slide 7: Current Trends
    slide = ppt_pres.slides.add_slide(ppt_pres.slide_layouts[1])
    slide.shapes.title.text = "Current Trends in Telecommunications"
    content = slide.placeholders[1]
    content.text = ""  # Clear default text
    p = content.text_frame.add_paragraph()
    set_font(p, font_size=Pt(20))
    p.text = "Current trends include:"
    
    current_trends = [
        "Remote Work: Increased demand for reliable connectivity as more employees work from home.",
        "Cybersecurity: Growing concerns over data security and privacy.",
        "Cloud Services: Shift towards cloud-based solutions for flexibility and scalability.",
        "Sustainability: Rising importance of environmentally friendly practices in telecom operations."
    ]
    
    for trend in current_trends:
        p = content.text_frame.add_paragraph()
        set_font(p, font_size=Pt(18))
        p.text = f"- {trend}"

    # Slide 8: Addressing Industry-Specific Problems
    slide = ppt_pres.slides.add_slide(ppt_pres.slide_layouts[1])
    slide.shapes.title.text = "Addressing Industry-Specific Problems"
    content = slide.placeholders[1]
    content.text = ""  # Clear default text
    p = content.text_frame.add_paragraph()
    set_font(p, font_size=Pt(20))
    p.text = "Accenture can contribute by:"
    
    solutions = [
        "Digital Transformation: Helping clients transition to digital-first operations.",
        "Cybersecurity Solutions: Providing robust security frameworks to protect client data.",
        "Sustainable Practices: Advising on green technologies and practices to reduce environmental impact.",
        "Innovation Hubs: Establishing centers for co-innovation to develop new solutions."
    ]
    
    for solution in solutions:
        p = content.text_frame.add_paragraph()
        set_font(p, font_size=Pt(18))
        p.text = f"- {solution}"

    # Slide 9: Conclusion
    slide = ppt_pres.slides.add_slide(ppt_pres.slide_layouts[1])
    slide.shapes.title.text = "Conclusion"
    content = slide.placeholders[1]
    content.text = ""  # Clear default text
    p = content.text_frame.add_paragraph()
    set_font(p, font_size=Pt(20))
    p.text = "In conclusion, understanding future trends in telecommunications is crucial for:"
    
    conclusions = [
        "Adapting project strategies to meet evolving client needs.",
        "Enhancing operational efficiency and service delivery.",
        "Building strong, future-ready partnerships with clients."
    ]
    
    for conclusion in conclusions:
        p = content.text_frame.add_paragraph()
        set_font(p, font_size=Pt(18))
        p.text = f"- {conclusion}"

    # Slide 10: Q&A
    slide = ppt_pres.slides.add_slide(ppt_pres.slide_layouts[1])
    slide.shapes.title.text = "Q&A"
    content = slide.placeholders[1]
    content.text = ""  # Clear default text
    p = content.text_frame.add_paragraph()
    set_font(p, font_size=Pt(20))
    p.text = "Thank you for your attention!"
    
    p = content.text_frame.add_paragraph()
    set_font(p, font_size=Pt(18))
    p.text = "Questions?"

    # Save the presentation
    ppt_pres.save(r'C:\Users\priti\OneDrive\Documents\Pic\Week3_OJT_yourname.pptx')

# Call the function to create the presentation
create_professional_presentation()